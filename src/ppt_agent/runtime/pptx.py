from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

from ppt_agent.domain.evidence import EvidenceItem, EvidencePack, FigureAsset, TableAsset
from ppt_agent.domain.models import Artifact, PptSpec, SlideSpec
from ppt_agent.runtime.styles import StylePreset, get_style

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def build_pptx(
    spec: PptSpec,
    output_path: Path,
    evidence_pack: EvidencePack | None = None,
    evidence_path: Path | None = None,
    debug_source_trace: bool = False,
) -> Artifact:
    output_path.parent.mkdir(parents=True, exist_ok=True)

    style = get_style(spec.theme)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    figures_by_id = _figures_by_id(evidence_pack)
    tables_by_id = _tables_by_id(evidence_pack)
    evidence_by_id = _evidence_by_id(evidence_pack)

    for slide_spec in spec.slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _paint_background(slide, style)
        _apply_accent(slide, style)
        layout = _resolve_layout(slide_spec)
        _render_layout(
            slide, slide_spec, layout, style,
            figures_by_id=figures_by_id,
            tables_by_id=tables_by_id,
            evidence_path=evidence_path,
        )
        _add_source_footer(slide, slide_spec, style, evidence_by_id=evidence_by_id)
        notes_text = _speaker_notes_with_source_trace(
            slide_spec,
            evidence_by_id=evidence_by_id,
            debug_source_trace=debug_source_trace,
        )
        if notes_text:
            slide.notes_slide.notes_text_frame.text = notes_text

    prs.save(output_path)
    return Artifact(path=output_path)


def _resolve_layout(slide_spec: SlideSpec) -> str:
    explicit_layout = (slide_spec.layout or "").strip()
    if explicit_layout:
        return explicit_layout
    if slide_spec.layout_hint:
        return slide_spec.layout_hint
    if slide_spec.content.figure_ids:
        return "figure_walkthrough"
    mapping = {
        "hero_image": "title_cover",
        "market_scene": "hero_image_plus_argument",
        "workspace_photo": "two_column_text_image",
        "customer_moment": "hero_image_plus_argument",
        "three_card_summary": "three_card_summary",
        "concept_explainer": "concept_explainer",
        "result_cards": "result_cards",
        "process_timeline": "process_timeline",
        "method_step_flow": "method_step_flow",
        "comparison_table": "comparison_table",
        "figure_walkthrough": "figure_walkthrough",
    }
    return mapping.get(slide_spec.visual_type, "two_column_text_image")


def _render_layout(
    slide,
    slide_spec: SlideSpec,
    layout: str,
    style: StylePreset,
    *,
    figures_by_id: dict[str, FigureAsset] | None = None,
    tables_by_id: dict[str, TableAsset] | None = None,
    evidence_path: Path | None = None,
) -> None:
    figure_layouts = {"figure_walkthrough", "figure_with_caption", "figure_caption", "two_column_figure", "method_figure_callouts"}
    if slide_spec.content.figure_ids and layout not in figure_layouts:
        layout = "figure_walkthrough"
    # Tables only get their own layout if no figure_ids are present (figures take priority).
    # When both figure_ids and table_ids exist, figures win; table data remains in slide_spec
    # for source footnotes and speaker notes but is not rendered as a dedicated table layout.
    if not slide_spec.content.figure_ids and slide_spec.content.table_ids and layout not in {"result_cards", "result_table_summary", "comparison_table"}:
        layout = "result_cards"
    renderer = {
        "hero": _render_academic_cover,
        "title_cover": _render_title_cover,
        "title-bullets": _render_title_bullets,
        "hero_image_plus_argument": _render_hero_image_plus_argument,
        "two_column_text_image": _render_two_column_text_image,
        # NOTE: 'concept_explainer' intentionally uses the three-card layout renderer.
        # The three-card summary is the closest generic layout for explaining concepts via
        # three key points.  A dedicated concept_explainer renderer can be added later if needed.
        "concept_explainer": _render_three_card_summary,
        "three_card_summary": _render_three_card_summary,
        "method_step_flow": _render_process_timeline,
        "process_timeline": _render_process_timeline,
        "comparison_table": _render_comparison_table,
        "figure_walkthrough": _render_method_figure_callouts,
        "figure_with_caption": _render_figure_with_caption,
        "figure_caption": _render_multi_figure_slide,
        "two_column_figure": _render_multi_figure_slide,
        "method_figure_callouts": _render_method_figure_callouts,
        "result_cards": _render_result_table_summary,
        "result_table_summary": _render_result_table_summary,
        "ablation_summary": _render_ablation_summary,
        "big_quote": _render_big_quote,
    }.get(layout)
    if renderer is None:
        _render_two_column_text_image(slide, slide_spec, style)
        return
    if layout in {"figure_walkthrough", "figure_with_caption", "figure_caption", "two_column_figure", "method_figure_callouts"}:
        renderer(slide, slide_spec, style, figures_by_id=figures_by_id or {}, evidence_path=evidence_path)
    elif layout in {"result_cards", "result_table_summary"}:
        renderer(slide, slide_spec, style, tables_by_id=tables_by_id or {}, evidence_path=evidence_path)
    else:
        renderer(slide, slide_spec, style)


def _render_academic_cover(slide, slide_spec: SlideSpec, style: StylePreset) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(4.45), SLIDE_HEIGHT)
    band.fill.solid()
    band.fill.fore_color.rgb = style.primary
    band.line.color.rgb = style.primary
    _add_textbox(slide, Inches(0.65), Inches(0.7), Inches(3.3), Inches(0.35), "PAPER BRIEF", 11, True, style.tertiary, font_name=style.font_name)
    _add_textbox(slide, Inches(0.65), Inches(1.35), Inches(3.35), Inches(3.0), slide_spec.title, 25, True, style.surface, font_name=style.font_name)
    _add_textbox(slide, Inches(0.65), Inches(4.65), Inches(3.25), Inches(1.1), _join_lines(_visible_bullets(slide_spec)), 13, False, style.surface, font_name=style.font_name)
    _add_textbox(slide, Inches(5.0), Inches(1.2), Inches(6.9), Inches(1.35), _slide_message(slide_spec), 24, True, style.primary, font_name=style.font_name)
    _draw_evidence_motif(slide, style, Inches(5.0), Inches(3.0), Inches(6.9), Inches(2.6))


def _render_title_cover(slide, slide_spec: SlideSpec, style: StylePreset) -> None:
    _add_textbox(slide, Inches(0.8), Inches(0.7), Inches(6.2), Inches(0.8), slide_spec.title, 28, True, style.surface, font_name=style.font_name)
    _add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(1.2), _slide_message(slide_spec), 18, False, style.surface, font_name=style.font_name)
    _add_textbox(slide, Inches(0.8), Inches(3.0), Inches(4.8), Inches(2.2), _join_lines(_visible_bullets(slide_spec)), 15, False, style.surface, font_name=style.font_name)
    _render_visual_area(slide, slide_spec, style, Inches(7.0), Inches(0.7), Inches(5.1), Inches(5.5), accent=style.tertiary)


def _render_title_bullets(slide, slide_spec: SlideSpec, style: StylePreset) -> None:
    _add_section_title(slide, slide_spec.title, style)
    _add_textbox(slide, Inches(0.85), Inches(1.15), Inches(11.4), Inches(0.7), _slide_message(slide_spec), 18, True, style.secondary, font_name=style.font_name)
    bullets = _visible_bullets(slide_spec)
    while len(bullets) < 3:
        bullets.append("")
    for index, bullet in enumerate(bullets[:3]):
        left = Inches(0.85 + index * 4.05)
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(2.15), Inches(3.45), Inches(3.2))
        card.fill.solid()
        card.fill.fore_color.rgb = style.surface
        card.line.color.rgb = style.tertiary if index == 1 else style.border
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.22), Inches(2.35), Inches(0.45), Inches(0.45))
        badge.fill.solid()
        badge.fill.fore_color.rgb = style.secondary if index != 1 else style.tertiary
        badge.line.color.rgb = badge.fill.fore_color.rgb
        _set_text(badge.text_frame, str(index + 1), 10, True, style.surface, center=True)
        _add_textbox(slide, left + Inches(0.25), Inches(3.0), Inches(2.95), Inches(1.7), bullet, 15, False, style.text_body, font_name=style.font_name)


def _render_hero_image_plus_argument(slide, slide_spec: SlideSpec, style: StylePreset) -> None:
    _add_section_title(slide, slide_spec.title, style)
    _add_textbox(slide, Inches(0.8), Inches(1.5), Inches(4.7), Inches(0.8), _slide_message(slide_spec), 22, True, style.primary, font_name=style.font_name)
    _add_textbox(slide, Inches(0.8), Inches(2.4), Inches(4.7), Inches(2.8), _join_lines(_visible_bullets(slide_spec)), 16, False, style.text_body, font_name=style.font_name)
    _render_visual_area(slide, slide_spec, style, Inches(6.0), Inches(1.3), Inches(6.1), Inches(4.9), accent=style.secondary)


def _render_two_column_text_image(slide, slide_spec: SlideSpec, style: StylePreset) -> None:
    _add_section_title(slide, slide_spec.title, style)
    _add_textbox(slide, Inches(0.8), Inches(1.4), Inches(5.1), Inches(0.7), _slide_message(slide_spec), 17, True, style.secondary, font_name=style.font_name)
    _add_textbox(slide, Inches(0.8), Inches(2.1), Inches(5.0), Inches(2.8), _join_lines(_visible_bullets(slide_spec)), 15, False, style.text_body, font_name=style.font_name)
    _render_visual_area(slide, slide_spec, style, Inches(6.2), Inches(1.4), Inches(5.8), Inches(4.8), accent=style.tertiary)


def _render_three_card_summary(slide, slide_spec: SlideSpec, style: StylePreset) -> None:
    _add_section_title(slide, slide_spec.title, style)
    _add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.4), Inches(0.6), _slide_message(slide_spec), 18, False, style.text_body, font_name=style.font_name)
    items = _visible_bullets(slide_spec)
    while len(items) < 3:
        items.append("")
    lefts = [Inches(0.8), Inches(4.45), Inches(8.1)]
    for idx, item in enumerate(items):
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, lefts[idx], Inches(2.0), Inches(3.1), Inches(3.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = style.surface
        shape.line.color.rgb = style.border
        _set_text(shape.text_frame, item, 16, False, style.primary)


def _render_process_timeline(slide, slide_spec: SlideSpec, style: StylePreset) -> None:
    _add_section_title(slide, slide_spec.title, style)
    _add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.2), Inches(0.7), _slide_message(slide_spec), 17, False, style.text_body, font_name=style.font_name)
    steps = _visible_bullets(slide_spec)
    while len(steps) < 3:
        steps.append("")
    start_left = Inches(1.0)
    top = Inches(3.5)
    width = Inches(2.6)
    for idx, step in enumerate(steps[:3]):
        left = start_left + Inches(3.0) * idx
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, Inches(2.4), Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = style.secondary
        circle.line.color.rgb = style.secondary
        _set_text(circle.text_frame, str(idx + 1), 12, True, style.surface, center=True)
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left - Inches(0.2), top, width, Inches(1.4))
        box.fill.solid()
        box.fill.fore_color.rgb = style.surface
        box.line.color.rgb = style.border
        _set_text(box.text_frame, step, 13, False, style.text_body)
        if idx < len(steps) - 1:
            connector = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left + Inches(0.8), Inches(2.55), Inches(1.6), Inches(0.3))
            connector.fill.solid()
            connector.fill.fore_color.rgb = style.tertiary
            connector.line.color.rgb = style.tertiary


def _render_comparison_table(slide, slide_spec: SlideSpec, style: StylePreset) -> None:
    _add_section_title(slide, slide_spec.title, style)
    _add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.2), Inches(0.7), _slide_message(slide_spec), 17, False, style.text_body, font_name=style.font_name)
    rows = 4
    table = slide.shapes.add_table(rows, 3, Inches(0.8), Inches(2.0), Inches(11.4), Inches(3.5)).table
    headers = ["Dimension", "Current State", "Target State"]
    for col, value in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = value
        cell.fill.solid()
        cell.fill.fore_color.rgb = style.primary
        _style_cell(cell, True, style.surface)
    points = _visible_bullets(slide_spec)
    rows_data = [
        ("Seller workflow", points[0] if len(points) > 0 else "Manual prep", slide_spec.bullets[0] if slide_spec.bullets else "AI-guided prep"),
        ("Manager visibility", points[1] if len(points) > 1 else "Lagging signals", slide_spec.bullets[1] if len(slide_spec.bullets) > 1 else "Weekly leading indicators"),
        ("Business impact", points[2] if len(points) > 2 else "", _slide_message(slide_spec)),
    ]
    for row_idx, row in enumerate(rows_data, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = style.surface if row_idx % 2 else style.background
            _style_cell(cell, False, style.text_body)


def _render_figure_with_caption(
    slide,
    slide_spec: SlideSpec,
    style: StylePreset,
    *,
    figures_by_id: dict[str, FigureAsset],
    evidence_path: Path | None,
) -> None:
    _add_section_title(slide, slide_spec.title, style)
    figure_id = slide_spec.content.figure_ids[0] if slide_spec.content.figure_ids else None
    figure = figures_by_id.get(figure_id or "")
    image_path = _resolve_evidence_asset_path(figure.path, evidence_path=evidence_path) if figure and figure.path else None
    bullets = _visible_bullets(slide_spec)
    has_bullets = bool(bullets)

    image_left = Inches(0.85)
    image_top = Inches(1.25)
    image_width = Inches(7.5 if has_bullets else 11.6)
    image_height = Inches(4.75)
    caption_top = Inches(6.1)

    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), image_left, image_top, width=image_width, height=image_height)
    else:
        _draw_missing_figure_placeholder(slide, slide_spec, style, image_left, image_top, image_width, image_height, figure_id=figure_id)

    caption = _figure_caption(figure, slide_spec, figure_id=figure_id)
    _add_textbox(slide, image_left, caption_top, image_width, Inches(0.55), caption, 10, False, style.text_muted, font_name=style.font_name)

    if has_bullets:
        message = _slide_message(slide_spec)
        _add_textbox(slide, Inches(8.65), Inches(1.35), Inches(3.65), Inches(0.9), message, 16, True, style.primary, font_name=style.font_name)
        _add_textbox(slide, Inches(8.65), Inches(2.3), Inches(3.65), Inches(3.0), _join_lines(bullets), 13, False, style.text_body, font_name=style.font_name)


def _render_multi_figure_slide(
    slide,
    slide_spec: SlideSpec,
    style: StylePreset,
    *,
    figures_by_id: dict[str, FigureAsset],
    evidence_path: Path | None,
) -> None:
    figure_ids = slide_spec.content.figure_ids
    if len(figure_ids) <= 1:
        _render_figure_with_caption(slide, slide_spec, style, figures_by_id=figures_by_id, evidence_path=evidence_path)
        return

    _add_section_title(slide, slide_spec.title, style)
    _add_textbox(slide, Inches(0.85), Inches(1.15), Inches(11.5), Inches(0.55), _slide_message(slide_spec), 14, True, style.primary, font_name=style.font_name)
    panels = [
        (Inches(0.85), Inches(1.9), Inches(5.65), Inches(3.65)),
        (Inches(6.85), Inches(1.9), Inches(5.65), Inches(3.65)),
    ]
    for figure_id, (left, top, width, height) in zip(figure_ids[:2], panels):
        figure = figures_by_id.get(figure_id)
        image_path = _resolve_evidence_asset_path(figure.path, evidence_path=evidence_path) if figure and figure.path else None
        frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
        frame.fill.solid()
        frame.fill.fore_color.rgb = style.surface
        frame.line.color.rgb = style.border
        if image_path and image_path.exists():
            _add_picture_contain(slide, image_path, left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), height - Inches(0.55))
        else:
            _draw_missing_figure_placeholder(slide, slide_spec, style, left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), height - Inches(0.55), figure_id=figure_id)
        caption = _figure_caption(figure, slide_spec, figure_id=figure_id)
        _add_textbox(slide, left + Inches(0.18), top + height - Inches(0.42), width - Inches(0.36), Inches(0.32), caption, 8, False, style.text_muted, font_name=style.font_name)
    bullets = _visible_bullets(slide_spec)
    if bullets:
        _add_textbox(slide, Inches(0.95), Inches(5.9), Inches(11.1), Inches(0.75), _join_lines(bullets), 12, False, style.text_body, font_name=style.font_name)


def _render_method_figure_callouts(
    slide,
    slide_spec: SlideSpec,
    style: StylePreset,
    *,
    figures_by_id: dict[str, FigureAsset],
    evidence_path: Path | None,
) -> None:
    _add_section_title(slide, slide_spec.title, style)
    figure_id = slide_spec.content.figure_ids[0] if slide_spec.content.figure_ids else None
    figure = figures_by_id.get(figure_id or "")
    image_path = _resolve_evidence_asset_path(figure.path, evidence_path=evidence_path) if figure and figure.path else None
    image_left, image_top, image_width, image_height = Inches(0.8), Inches(1.25), Inches(7.3), Inches(4.8)
    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), image_left, image_top, width=image_width, height=image_height)
    else:
        _draw_missing_figure_placeholder(slide, slide_spec, style, image_left, image_top, image_width, image_height, figure_id=figure_id)

    callouts = list(slide_spec.content.callouts)
    if not callouts:
        callouts = [
            {"label": str(index + 1), "text": bullet}
            for index, bullet in enumerate(_visible_bullets(slide_spec))
        ]
    _add_textbox(slide, Inches(8.45), Inches(1.25), Inches(3.9), Inches(0.7), _slide_message(slide_spec), 15, True, style.primary, font_name=style.font_name)
    for index, callout in enumerate(callouts[:3]):
        # Handle both VisualCallout objects and plain dicts (fallback case)
        if isinstance(callout, dict):
            label = callout.get("label", "") or str(index + 1)
            text = callout.get("text", "")
        else:
            label = getattr(callout, "label", "") or str(index + 1)
            text = getattr(callout, "text", "") or ""
        top = Inches(2.05 + index * 1.2)
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(8.5), top, Inches(0.42), Inches(0.42))
        badge.fill.solid()
        badge.fill.fore_color.rgb = style.tertiary
        badge.line.color.rgb = style.tertiary
        _set_text(badge.text_frame, str(label), 10, True, style.surface, center=True)
        _add_textbox(slide, Inches(9.05), top - Inches(0.05), Inches(3.1), Inches(0.8), text, 12, False, style.text_body, font_name=style.font_name)
    caption = _figure_caption(figure, slide_spec, figure_id=figure_id)
    reason = slide_spec.content.visual_reason or slide_spec.image_rationale
    footer = f"{caption}\n{reason}" if reason else caption
    _add_textbox(slide, image_left, Inches(6.15), Inches(11.4), Inches(0.6), footer, 10, False, style.text_muted, font_name=style.font_name)


def _render_result_table_summary(
    slide, slide_spec: SlideSpec, style: StylePreset,
    *, tables_by_id: dict[str, TableAsset] | None = None, evidence_path: Path | None = None,
) -> None:
    _add_section_title(slide, slide_spec.title, style)
    _add_textbox(slide, Inches(0.8), Inches(1.25), Inches(11.4), Inches(0.7), _slide_message(slide_spec), 18, True, style.primary, font_name=style.font_name)

    # Try to render table images from evidence
    table_rendered = False
    if tables_by_id and slide_spec.content.table_ids:
        for tid in slide_spec.content.table_ids[:1]:
            table = tables_by_id.get(tid)
            if table and table.path:
                image_path = _resolve_evidence_asset_path(table.path, evidence_path=evidence_path)
                if image_path and image_path.exists():
                    try:
                        slide.shapes.add_picture(str(image_path), Inches(0.8), Inches(2.1), width=Inches(11.4), height=Inches(4.0))
                        table_rendered = True
                    except Exception:
                        pass  # fall through to text summary

    if not table_rendered:
        summaries = [
            _result_summary_text(item)
            for item in slide_spec.content.result_summary
        ] or [
            str(metric.get("finding") or metric.get("value") or metric.get("metric") or metric)
            for metric in slide_spec.content.metrics
        ] or _visible_bullets(slide_spec)
        while len(summaries) < 3:
            summaries.append("Result detail not provided by evidence")
        lefts = [Inches(0.85), Inches(4.55), Inches(8.25)]
        for index, summary in enumerate(summaries[:3]):
            card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, lefts[index], Inches(2.25), Inches(3.25), Inches(2.55))
            card.fill.solid()
            card.fill.fore_color.rgb = style.surface
            card.line.color.rgb = style.border
            _set_text(card.text_frame, summary, 15, True if index == 0 else False, style.primary if index == 0 else style.text_body)
    trace = "Table summary: " + ", ".join(slide_spec.content.table_ids) if slide_spec.content.table_ids else "Table summary"
    if slide_spec.content.visual_reason:
        trace = f"{trace} - {slide_spec.content.visual_reason}"
    citations = "Citations: " + ", ".join(citation.evidence_id for citation in slide_spec.citations[:4])
    _add_textbox(slide, Inches(0.85), Inches(5.25), Inches(11.3), Inches(0.45), trace, 11, False, style.text_muted, font_name=style.font_name)
    _add_textbox(slide, Inches(0.85), Inches(5.75), Inches(11.3), Inches(0.45), citations, 10, False, style.text_muted, font_name=style.font_name)


def _render_ablation_summary(slide, slide_spec: SlideSpec, style: StylePreset) -> None:
    _add_section_title(slide, slide_spec.title, style)
    _add_textbox(slide, Inches(0.8), Inches(1.25), Inches(11.4), Inches(0.65), _slide_message(slide_spec), 17, True, style.primary, font_name=style.font_name)
    points = _visible_bullets(slide_spec)
    results = [_result_summary_text(item) for item in slide_spec.content.result_summary]
    columns = [
        ("What changed", points[0] if len(points) > 0 else "Variant not provided by evidence"),
        ("What happened", results[0] if results else (points[1] if len(points) > 1 else "Observation not provided by evidence")),
        ("Interpretation", points[2] if len(points) > 2 else slide_spec.content.visual_reason or "Interpretation requires evidence review"),
    ]
    for index, (heading, text) in enumerate(columns):
        left = Inches(0.85 + index * 4.0)
        panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(2.15), Inches(3.45), Inches(3.35))
        panel.fill.solid()
        panel.fill.fore_color.rgb = style.surface
        panel.line.color.rgb = style.tertiary if index == 1 else style.border
        _set_text(panel.text_frame, f"{heading}\n{text}", 14, True, style.text_body)
    _add_textbox(slide, Inches(0.85), Inches(5.8), Inches(11.2), Inches(0.5), ", ".join([c.evidence_id for c in slide_spec.citations[:4]]), 10, False, style.text_muted, font_name=style.font_name)


def _render_big_quote(slide, slide_spec: SlideSpec, style: StylePreset) -> None:
    _add_textbox(slide, Inches(0.9), Inches(0.65), Inches(10.8), Inches(0.45), "KEY TAKEAWAY", 11, True, style.tertiary, font_name=style.font_name)
    _add_textbox(slide, Inches(0.9), Inches(1.45), Inches(10.9), Inches(2.0), _slide_message(slide_spec), 26, True, style.primary, font_name=style.font_name)
    bullets = _visible_bullets(slide_spec)
    lefts = [Inches(0.95), Inches(4.55), Inches(8.15)]
    for index, bullet in enumerate((bullets + ["", "", ""])[:3]):
        panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, lefts[index], Inches(4.0), Inches(3.2), Inches(1.45))
        panel.fill.solid()
        panel.fill.fore_color.rgb = style.surface
        panel.line.color.rgb = style.tertiary if index == 0 else style.border
        _set_text(panel.text_frame, bullet, 13, False, style.text_body)


def _result_summary_text(item) -> str:
    metric = getattr(item, "metric", None) if not isinstance(item, dict) else item.get("metric")
    finding = getattr(item, "finding", None) if not isinstance(item, dict) else item.get("finding")
    if metric and finding:
        return f"{metric}: {finding}"
    return finding or metric or str(item)


def _paint_background(slide, style: StylePreset) -> None:
    background = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    background.fill.solid()
    background.fill.fore_color.rgb = style.background
    background.line.color.rgb = style.background
    slide.shapes._spTree.remove(background._element)
    slide.shapes._spTree.insert(2, background._element)


def _apply_accent(slide, style: StylePreset) -> None:
    pattern = style.accent_pattern
    if pattern == "gradient_band":
        band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.15))
        band.fill.solid()
        band.fill.fore_color.rgb = style.secondary
        band.line.color.rgb = style.secondary
        slide.shapes._spTree.remove(band._element)
        slide.shapes._spTree.insert(2, band._element)
    elif pattern == "side_bar":
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(0.15), SLIDE_HEIGHT)
        bar.fill.solid()
        bar.fill.fore_color.rgb = style.primary
        bar.line.color.rgb = style.primary
        slide.shapes._spTree.remove(bar._element)
        slide.shapes._spTree.insert(2, bar._element)
    elif pattern == "corner_stripe":
        stripe = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            SLIDE_WIDTH - Inches(2.5), 0,
            Inches(2.5), Inches(0.15),
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = style.secondary
        stripe.line.color.rgb = style.secondary
        slide.shapes._spTree.remove(stripe._element)
        slide.shapes._spTree.insert(2, stripe._element)
    elif pattern == "underline":
        line = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0.8), Inches(1.05),
            Inches(11.4), Inches(0.02),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = style.secondary
        line.line.color.rgb = style.secondary


def _add_section_title(slide, title: str, style: StylePreset) -> None:
    _add_textbox(slide, Inches(0.8), Inches(0.45), Inches(11.4), Inches(0.6), title, 24, True, style.primary, font_name=style.font_name)


def _render_visual_area(slide, slide_spec: SlideSpec, style: StylePreset, left, top, width, height, *, accent: RGBColor) -> None:
    local_path = (slide_spec.resolved_asset or {}).get("local_path")
    if local_path and Path(local_path).exists():
        _add_picture(slide, style, local_path, left, top, width, height, caption=slide_spec.image_caption or slide_spec.core_message)
        return
    _draw_visual_placeholder(slide, slide_spec, style, left, top, width, height, accent=accent)


def _draw_visual_placeholder(slide, slide_spec: SlideSpec, style: StylePreset, left, top, width, height, *, accent: RGBColor) -> None:
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    frame.fill.solid()
    frame.fill.fore_color.rgb = style.surface
    frame.line.color.rgb = accent

    layout = _resolve_layout(slide_spec)
    if layout == "title_cover":
        _draw_cover_fallback(slide, slide_spec, style, left, top, width, height, accent=accent)
    elif layout == "hero_image_plus_argument":
        _draw_hero_fallback(slide, slide_spec, style, left, top, width, height, accent=accent)
    else:
        _draw_two_column_fallback(slide, slide_spec, style, left, top, width, height, accent=accent)


def _draw_cover_fallback(slide, slide_spec: SlideSpec, style: StylePreset, left, top, width, height, *, accent: RGBColor) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(1.0))
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.color.rgb = accent

    circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.4), top + Inches(1.35), Inches(1.2), Inches(1.2))
    circle.fill.solid()
    circle.fill.fore_color.rgb = style.background
    circle.line.color.rgb = style.border

    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left + Inches(1.9),
        top + Inches(1.2),
        width - Inches(2.4),
        height - Inches(1.7),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = style.background
    card.line.color.rgb = style.border
    label = slide_spec.image_caption or _slide_message(slide_spec) or slide_spec.title
    _set_text(card.text_frame, label, 20, True, style.primary, font_name=style.font_name)


def _draw_hero_fallback(slide, slide_spec: SlideSpec, style: StylePreset, left, top, width, height, *, accent: RGBColor) -> None:
    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + Inches(0.3), top + Inches(0.3), width - Inches(0.6), height - Inches(0.6))
    hero.fill.solid()
    hero.fill.fore_color.rgb = style.background
    hero.line.color.rgb = style.border

    for idx in range(3):
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.CHEVRON,
            left + Inches(0.6) + Inches(1.15) * idx,
            top + Inches(0.7),
            Inches(0.9),
            Inches(0.32),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent if idx % 2 == 0 else style.tertiary
        bar.line.color.rgb = bar.fill.fore_color.rgb

    message = slide_spec.image_caption or _slide_message(slide_spec) or slide_spec.title
    textbox = slide.shapes.add_textbox(left + Inches(0.7), top + Inches(1.5), width - Inches(1.4), Inches(2.2))
    _set_text(textbox.text_frame, message, 18, True, style.primary, font_name=style.font_name)


def _draw_two_column_fallback(slide, slide_spec: SlideSpec, style: StylePreset, left, top, width, height, *, accent: RGBColor) -> None:
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + Inches(0.3), top + Inches(0.3), width - Inches(0.6), height - Inches(0.6))
    panel.fill.solid()
    panel.fill.fore_color.rgb = style.background
    panel.line.color.rgb = style.border

    role = (slide_spec.role or slide_spec.visual_type or "Evidence").replace("_", " ").title()
    role_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + Inches(0.7), top + Inches(0.8), Inches(2.4), Inches(0.55))
    role_box.fill.solid()
    role_box.fill.fore_color.rgb = style.surface
    role_box.line.color.rgb = accent
    _set_text(role_box.text_frame, role, 13, True, accent, center=True, font_name=style.font_name)

    message = slide_spec.image_caption or _slide_message(slide_spec) or slide_spec.title
    caption = slide.shapes.add_textbox(left + Inches(0.7), top + Inches(3.6), width - Inches(1.4), Inches(1.0))
    _set_text(caption.text_frame, message, 16, False, style.text_body, center=True, font_name=style.font_name)


def _draw_evidence_motif(slide, style: StylePreset, left, top, width, height) -> None:
    for index in range(4):
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left + Inches(0.15),
            top + Inches(0.25 + index * 0.52),
            width - Inches(0.3 + index * 0.45),
            Inches(0.22),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = style.tertiary if index in {0, 2} else style.secondary
        bar.line.color.rgb = bar.fill.fore_color.rgb
    node_positions = [(0.6, 1.95), (2.4, 1.25), (4.2, 2.05), (5.9, 1.45)]
    for x, y in node_positions:
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(x), top + Inches(y), Inches(0.55), Inches(0.55))
        node.fill.solid()
        node.fill.fore_color.rgb = style.surface
        node.line.color.rgb = style.tertiary


def _draw_missing_figure_placeholder(slide, slide_spec: SlideSpec, style: StylePreset, left, top, width, height, *, figure_id: str | None) -> None:
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    frame.fill.solid()
    frame.fill.fore_color.rgb = style.surface
    frame.line.color.rgb = style.tertiary
    message = f"Figure image missing: {figure_id or 'no figure_id'}"
    if slide_spec.image_caption:
        message = f"{message}\n{slide_spec.image_caption}"
    _set_text(frame.text_frame, message, 18, True, style.primary, center=True, font_name=style.font_name)


def _resolve_evidence_asset_path(path: str, *, evidence_path: Path | None) -> Path | None:
    candidate = Path(path)
    if candidate.is_absolute():
        return candidate if candidate.exists() else None
    if evidence_path is None:
        return candidate if candidate.exists() else None
    if candidate.exists():
        return candidate
    resolved = evidence_path.parent / candidate
    if resolved.exists():
        return resolved
    # Fallback: try resolving against the ingest directory
    # evidence_path is typically: workspace/.ppt-agent/data/evidence/<source_id>/evidence.json
    # ingest images are typically: workspace/.ppt-agent/ingest/<source_id>/images/
    try:
        evidence_dir = evidence_path.parent  # .../evidence/<source_id>/
        source_id = evidence_dir.name
        workspace = evidence_dir.parent.parent.parent  # workspace/
        ingest_dir = workspace / ".ppt-agent" / "ingest" / source_id
        for sub in ("images", "assets", ""):
            ingest_candidate = (ingest_dir / sub / candidate.name) if sub else (ingest_dir / candidate.name)
            if ingest_candidate.exists():
                return ingest_candidate
        # Also try the candidate relative to ingest_dir
        ingest_candidate = ingest_dir / candidate
        if ingest_candidate.exists():
            return ingest_candidate
    except (ValueError, OSError):
        pass
    return None


def _figures_by_id(evidence_pack: EvidencePack | None) -> dict[str, FigureAsset]:
    if evidence_pack is None:
        return {}
    return {figure.id: figure for figure in evidence_pack.figures}


def _tables_by_id(evidence_pack: EvidencePack | None) -> dict[str, TableAsset]:
    if evidence_pack is None:
        return {}
    return {table.id: table for table in evidence_pack.tables}


def _evidence_by_id(evidence_pack: EvidencePack | None) -> dict[str, EvidenceItem]:
    if evidence_pack is None:
        return {}
    return {item.id: item for item in evidence_pack.evidence_items()}


def _speaker_notes_with_source_trace(
    slide_spec: SlideSpec,
    *,
    evidence_by_id: dict[str, EvidenceItem],
    debug_source_trace: bool,
) -> str:
    lines = _source_trace_lines(slide_spec, evidence_by_id=evidence_by_id)
    parts = []
    if slide_spec.speaker_notes.strip():
        parts.append(slide_spec.speaker_notes.strip())
    if slide_spec.supporting_points:
        parts.append("Presenter details:\n" + "\n".join(f"- {point}" for point in slide_spec.supporting_points[:5]))
    if debug_source_trace and lines:
        parts.append("Source Trace:\n" + "\n".join(lines))
    return "\n\n".join(parts)


def _add_source_footer(slide, slide_spec: SlideSpec, style: StylePreset, *, evidence_by_id: dict[str, EvidenceItem]) -> None:
    footer = _compact_source_footer(slide_spec, evidence_by_id=evidence_by_id)
    if footer:
        _add_textbox(slide, Inches(0.8), Inches(6.95), Inches(11.6), Inches(0.25), footer, 8, False, style.text_muted, font_name=style.font_name)


def _compact_source_footer(slide_spec: SlideSpec, *, evidence_by_id: dict[str, EvidenceItem]) -> str:
    by_source: dict[str, set[int]] = {}
    for citation in slide_spec.citations:
        item = evidence_by_id.get(citation.evidence_id)
        source = citation.source_file or (item.source_file if item is not None else None)
        page = citation.page if citation.page is not None else (item.page if item is not None else None)
        if not source or page is None:
            continue
        by_source.setdefault(Path(source).name, set()).add(page)
    if not by_source:
        return ""
    parts = []
    for source, pages in by_source.items():
        page_text = ", ".join(f"p.{page}" for page in sorted(pages))
        parts.append(f"{source} {page_text}")
    return "Source: " + "; ".join(parts[:2])


def _slide_message(slide_spec: SlideSpec) -> str:
    return slide_spec.message or slide_spec.core_message


def _visible_bullets(slide_spec: SlideSpec) -> list[str]:
    bullets = slide_spec.content.bullets or slide_spec.bullets
    return [bullet for bullet in bullets if bullet][:3]


def _source_trace_lines(slide_spec: SlideSpec, *, evidence_by_id: dict[str, EvidenceItem]) -> list[str]:
    lines: list[str] = []
    seen: set[str] = set()
    for citation in slide_spec.citations:
        line = _format_source_trace(citation.evidence_id, source_file=citation.source_file, page=citation.page, evidence_by_id=evidence_by_id)
        if line not in seen:
            lines.append(line)
            seen.add(line)
    return lines


def _format_source_trace(
    evidence_id: str,
    *,
    source_file: str | None,
    page: int | None,
    evidence_by_id: dict[str, EvidenceItem],
) -> str:
    item = evidence_by_id.get(evidence_id)
    if item is None:
        if source_file:
            page_text = f" p.{page}" if page is not None else ""
            return f"Source: {source_file}{page_text} {evidence_id} unresolved"
        return f"Source: unresolved {evidence_id}"

    resolved_source = source_file or item.source_file
    resolved_page = page if page is not None else item.page
    page_text = f" p.{resolved_page}" if resolved_page is not None else ""
    caption = _evidence_caption(item)
    caption_text = f" - {caption}" if caption else ""
    return f"Source: {resolved_source}{page_text} {evidence_id}{caption_text}"


def _evidence_caption(item: EvidenceItem) -> str:
    for attr in ("caption", "heading", "text"):
        value = getattr(item, attr, None)
        if value:
            text = " ".join(str(value).split())
            return text[:120]
    return ""


def _figure_caption(figure: FigureAsset | None, slide_spec: SlideSpec, *, figure_id: str | None) -> str:
    if figure is None:
        return f"Source: missing figure evidence for {figure_id or 'unspecified figure'}"
    caption = figure.caption or slide_spec.image_caption or figure.text or figure.id
    source = figure.source_file
    page = f", p. {figure.page}" if figure.page is not None else ""
    return f"{caption} | Source: {source}{page}"


def _add_picture(slide, style: StylePreset, image_path: str, left, top, width, height, *, caption: str = "") -> None:
    slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    if caption:
        caption_box = slide.shapes.add_textbox(left, top + height - Inches(0.45), width, Inches(0.4))
        fill = caption_box.fill
        fill.solid()
        fill.fore_color.rgb = style.primary
        _set_text(caption_box.text_frame, caption, 11, False, style.surface, center=True, font_name=style.font_name)


def _add_picture_contain(slide, image_path: Path, left, top, width, height) -> None:
    try:
        from PIL import Image

        with Image.open(image_path) as image:
            image_width, image_height = image.size
    except Exception:
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
        return
    if not image_width or not image_height:
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
        return
    box_ratio = width / height
    image_ratio = image_width / image_height
    if image_ratio >= box_ratio:
        draw_width = width
        draw_height = width / image_ratio
    else:
        draw_height = height
        draw_width = height * image_ratio
    draw_left = left + (width - draw_width) / 2
    draw_top = top + (height - draw_height) / 2
    slide.shapes.add_picture(str(image_path), int(draw_left), int(draw_top), width=int(draw_width), height=int(draw_height))


def _add_textbox(slide, left, top, width, height, text: str, size: int, bold: bool, color: RGBColor, *, font_name: str = "Calibri"):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    _set_text(textbox.text_frame, text, size, bold, color, font_name=font_name)
    return textbox


def _set_text(text_frame, text: str, size: int, bold: bool, color: RGBColor, *, center: bool = False, font_name: str = "Calibri") -> None:
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text or " "
    if center:
        paragraph.alignment = 1
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def _style_cell(cell, bold: bool, color: RGBColor) -> None:
    paragraph = cell.text_frame.paragraphs[0]
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    run.font.bold = bold
    run.font.size = Pt(12)
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _join_lines(items: list[str]) -> str:
    return "\n".join(f"- {item}" for item in items if item)
