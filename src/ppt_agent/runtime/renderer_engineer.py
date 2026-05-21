from __future__ import annotations

import re
from pathlib import Path

import httpx
from pydantic import BaseModel, Field, ValidationError

from ppt_agent.llm.planner import PlannerConfigError
from ppt_agent.runtime.agent_llm import AgentLlmConfig, generate_agent_json
from ppt_agent.runtime.agent_skills import AgentSkillAssignments, skill_context_for_agent


class RendererEngineerArtifact(BaseModel):
    agent: str = "renderer_engineer"
    model: str = "deepseek-v4-pro"
    code_context: list[dict] = Field(default_factory=list)
    supported_layouts: list[str] = Field(default_factory=list)
    required_layouts: list[str] = Field(default_factory=list)
    gaps: list[dict] = Field(default_factory=list)
    extension_plan: list[dict] = Field(default_factory=list)
    generated_scripts: list[dict] = Field(default_factory=list)
    page_generator_contract: list[str] = Field(default_factory=list)
    risk_level: str = "low"


def renderer_engineer_agent_with_llm(
    slides_ir: dict,
    page_design: dict,
    config: AgentLlmConfig,
    skill_assignments: AgentSkillAssignments,
) -> dict:
    code_context = renderer_code_context()
    llm_payload = _llm_or_none(
        "renderer_engineer",
        RendererEngineerArtifact,
        system_prompt=renderer_engineer_prompt(),
        user_payload={
            "slides_ir_summary": slides_ir_summary(slides_ir),
            "page_design": page_design,
            "renderer_code_context": code_context,
            "assigned_skills": skill_context_for_agent(skill_assignments, "renderer_engineer"),
        },
        config=config,
    )
    if llm_payload is not None:
        return llm_payload
    return renderer_engineer_agent(slides_ir, page_design, code_context=code_context)


def renderer_engineer_agent(slides_ir: dict, page_design: dict, *, code_context: list[dict] | None = None) -> dict:
    supported = supported_renderer_layouts()
    required = sorted(
        {
            slide.get("layout")
            for slide in page_design.get("slides") or []
            if isinstance(slide, dict) and slide.get("layout")
        }
    )
    gaps = [
        {
            "layout": layout,
            "severity": "warning",
            "message": "Page Designer selected a layout that is not explicitly supported by the current PPTX renderer.",
            "suggested_file": "src/ppt_agent/runtime/pptx.py",
        }
        for layout in required
        if layout not in supported
    ]
    extension_plan = [
        {
            "kind": "layout_renderer",
            "target_file": "src/ppt_agent/runtime/pptx.py",
            "function": f"_render_{str(gap['layout']).replace('-', '_')}",
            "reason": gap["message"],
        }
        for gap in gaps
    ]
    generated_scripts = [_layout_inspection_script()] if gaps else []
    return {
        "agent": "renderer_engineer",
        "model": "deepseek-v4-pro",
        "code_context": code_context or [],
        "supported_layouts": supported,
        "required_layouts": required,
        "gaps": gaps,
        "extension_plan": extension_plan,
        "generated_scripts": generated_scripts,
        "page_generator_contract": [
            "Page Generator must remain deterministic and must not call an LLM.",
            "Page Generator should consume slides_ir.json plus page_design.json only.",
            "Renderer changes should be isolated to runtime renderer helpers or explicit support scripts.",
            "Unsupported layouts should degrade to a nonblank fallback rather than creating blank slides.",
        ],
        "risk_level": "medium" if gaps else "low",
    }


def renderer_engineer_prompt() -> str:
    return (
        "You are the Renderer Engineer Agent and should use deepseek-v4-pro when available. "
        "You may inspect Page Generator and PPTX renderer code context, but you must not rewrite slide content. "
        "Return JSON matching {agent, model, code_context, supported_layouts, required_layouts, gaps, "
        "extension_plan, generated_scripts, page_generator_contract, risk_level}. "
        "Identify whether the current deterministic renderer can implement page_design.json. "
        "If extra renderer scripts or layout functions are needed, describe source-code changes in extension_plan and "
        "put task-local helper scripts in generated_scripts with path, purpose, status, and optional content. "
        "Do not mutate slides_ir, page_design, PptSpec, PPTX, or repository source files. "
        "Keep recommendations concrete enough for a code agent to implement later."
    )


def validate_renderer_engineer_report(renderer_report: dict, page_design: dict) -> list[str]:
    issues = _model_validation_issues(RendererEngineerArtifact, renderer_report, "renderer_engineer_report")
    required = {
        slide.get("layout")
        for slide in page_design.get("slides") or []
        if isinstance(slide, dict) and slide.get("layout")
    }
    reported_required = set(renderer_report.get("required_layouts") or [])
    if required and not required.issubset(reported_required):
        issues.append("renderer_engineer_report missing one or more required page_design layouts")
    if renderer_report.get("risk_level") not in {"low", "medium", "high"}:
        issues.append("renderer_engineer_report risk_level must be low, medium, or high")
    return issues


def write_renderer_scripts(renderer_report: dict, scripts_dir: Path) -> None:
    scripts = renderer_report.get("generated_scripts") or []
    if not scripts:
        return
    scripts_dir.mkdir(parents=True, exist_ok=True)
    for index, script in enumerate(scripts, start=1):
        if not isinstance(script, dict) or not script.get("content"):
            continue
        raw_name = Path(str(script.get("path") or f"renderer_helper_{index}.py")).name
        safe_name = re.sub(r"[^A-Za-z0-9_.-]", "_", raw_name) or f"renderer_helper_{index}.py"
        if not safe_name.endswith(".py"):
            safe_name = f"{safe_name}.py"
        (scripts_dir / safe_name).write_text(str(script["content"]), encoding="utf-8")


def supported_renderer_layouts() -> list[str]:
    return [
        "hero",
        "title_cover",
        "title-bullets",
        "hero_image_plus_argument",
        "two_column_text_image",
        "concept_explainer",
        "three_card_summary",
        "method_step_flow",
        "process_timeline",
        "comparison_table",
        "figure_walkthrough",
        "figure_with_caption",
        "figure_caption",
        "two_column_figure",
        "method_figure_callouts",
        "result_cards",
        "result_table_summary",
        "ablation_summary",
        "big_quote",
    ]


def renderer_code_context() -> list[dict]:
    files = [
        Path(__file__).with_name("pptx.py"),
        Path(__file__).with_name("multi_agent_pipeline.py"),
    ]
    context = []
    for file_path in files:
        try:
            text = file_path.read_text(encoding="utf-8", errors="replace")
        except OSError:
            continue
        snippets = []
        for marker in ("def _render_layout", "def _resolve_layout", "def _page_generator_agent"):
            index = text.find(marker)
            if index >= 0:
                snippets.append(text[index : index + 2200])
        context.append(
            {
                "path": str(file_path),
                "purpose": "renderer/page-generator code visible to Renderer Engineer",
                "snippets": snippets[:3],
            }
        )
    return context


def slides_ir_summary(slides_ir: dict) -> dict:
    slides = slides_ir.get("slides") or []
    return {
        "deck": slides_ir.get("deck") or {},
        "theme": slides_ir.get("theme") or {},
        "slide_count": len(slides),
        "layouts": sorted({slide.get("layout") for slide in slides if slide.get("layout")}),
        "visual_types": sorted({slide.get("visual_type") for slide in slides if slide.get("visual_type")}),
    }


def _layout_inspection_script() -> dict[str, str]:
    return {
        "path": "scripts/inspect_pptx_layouts.py",
        "purpose": "Optional helper to inspect rendered slide layouts and flag blank or unsupported layout output.",
        "status": "proposed",
        "content": (
            "from __future__ import annotations\n\n"
            "from pathlib import Path\n"
            "from pptx import Presentation\n\n"
            "def inspect_pptx(path: str) -> dict:\n"
            "    prs = Presentation(path)\n"
            "    return {\n"
            "        'path': str(Path(path)),\n"
            "        'slide_count': len(prs.slides),\n"
            "        'shapes_per_slide': [len(slide.shapes) for slide in prs.slides],\n"
            "    }\n\n"
            "if __name__ == '__main__':\n"
            "    import json, sys\n"
            "    print(json.dumps(inspect_pptx(sys.argv[1]), ensure_ascii=False, indent=2))\n"
        ),
    }


def _llm_or_none(
    agent_name: str,
    model: type[BaseModel],
    *,
    system_prompt: str,
    user_payload: dict,
    config: AgentLlmConfig,
) -> dict | None:
    try:
        payload = generate_agent_json(
            agent_name,
            system_prompt=system_prompt,
            user_payload=user_payload,
            config=config,
        )
        if payload is None:
            return None
        return model.model_validate(payload).model_dump(mode="json")
    except (httpx.HTTPError, PlannerConfigError, ValueError, ValidationError):
        if config.fallback_to_deterministic:
            return None
        raise


def _model_validation_issues(model: type[BaseModel], payload: dict, label: str) -> list[str]:
    try:
        model.model_validate(payload)
    except ValidationError as exc:
        return [f"{label} schema error: {error['loc']} {error['msg']}" for error in exc.errors()]
    return []
