"""图表生成器 - 用于生成PPT中的图表"""

from typing import Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData


# 汇报专用配色方案
CHART_COLORS = [
    RGBColor(0x1A, 0x36, 0x5D),  # #1a365d 深蓝
    RGBColor(0x2B, 0x6C, 0xB0),  # #2b6cb0 浅蓝
    RGBColor(0x31, 0x82, 0xCE),  # #3182ce 强调蓝
    RGBColor(0x38, 0xA1, 0x69),  # #38a169 绿色
    RGBColor(0xD6, 0x9E, 0x2E),  # #d69e2e 黄色
    RGBColor(0xE5, 0x3E, 0x3E),  # #e53e3e 红色
    RGBColor(0x80, 0x5A, 0xD5),  # #805ad5 紫色
    RGBColor(0xDD, 0x6B, 0x20),  # #dd6b20 橙色
]


class ChartGenerator:
    """图表生成器"""
    
    def __init__(self):
        self.colors = CHART_COLORS
    
    def generate_bar_chart(self, slide: Any, data: dict, 
                           left: float = 1.5, top: float = 2,
                           width: float = 10, height: float = 4.5) -> Any:
        """生成柱状图
        
        Args:
            slide: 幻灯片对象
            data: 数据字典 {"categories": [...], "series": {"name": [...]}}
            left, top, width, height: 位置和大小
            
        Returns:
            图表对象
        """
        chart_data = CategoryChartData()
        chart_data.categories = data.get("categories", [])
        
        series = data.get("series", {})
        for name, values in series.items():
            chart_data.add_series(name, values)
        
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(left), Inches(top), Inches(width), Inches(height),
            chart_data
        ).chart
        
        # 设置样式
        self._apply_chart_style(chart)
        
        return chart
    
    def generate_line_chart(self, slide: Any, data: dict,
                            left: float = 1.5, top: float = 2,
                            width: float = 10, height: float = 4.5) -> Any:
        """生成折线图
        
        Args:
            slide: 幻灯片对象
            data: 数据字典 {"categories": [...], "series": {"name": [...]}}
            left, top, width, height: 位置和大小
            
        Returns:
            图表对象
        """
        chart_data = CategoryChartData()
        chart_data.categories = data.get("categories", [])
        
        series = data.get("series", {})
        for name, values in series.items():
            chart_data.add_series(name, values)
        
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS,
            Inches(left), Inches(top), Inches(width), Inches(height),
            chart_data
        ).chart
        
        # 设置样式
        self._apply_chart_style(chart)
        
        return chart
    
    def generate_pie_chart(self, slide: Any, data: dict,
                           left: float = 3, top: float = 2,
                           width: float = 7, height: float = 5) -> Any:
        """生成饼图
        
        Args:
            slide: 幻灯片对象
            data: 数据字典 {"categories": [...], "values": [...]}
            left, top, width, height: 位置和大小
            
        Returns:
            图表对象
        """
        chart_data = CategoryChartData()
        chart_data.categories = data.get("categories", [])
        chart_data.add_series("占比", data.get("values", []))
        
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE,
            Inches(left), Inches(top), Inches(width), Inches(height),
            chart_data
        ).chart
        
        # 设置样式
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        
        # 设置颜色
        plot = chart.plots[0]
        for i, point in enumerate(plot.series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = self.colors[i % len(self.colors)]
        
        return chart
    
    def generate_table(self, slide: Any, data: dict,
                       left: float = 1.5, top: float = 2,
                       width: float = 10, height: float = 4.5) -> Any:
        """生成表格
        
        Args:
            slide: 幻灯片对象
            data: 数据字典 {"headers": [...], "rows": [[...], ...]}
            left, top, width, height: 位置和大小
            
        Returns:
            表格对象
        """
        headers = data.get("headers", [])
        rows = data.get("rows", [])
        
        # 计算行数和列数
        num_rows = len(rows) + 1  # +1 for header
        num_cols = len(headers) if headers else len(rows[0]) if rows else 0
        
        # 添加表格
        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        table = table_shape.table
        
        # 设置表头
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = CHART_COLORS[0]
            
            # 设置字体
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    run.font.bold = True
                    run.font.size = Pt(14)
        
        # 设置数据行
        for i, row in enumerate(rows):
            for j, value in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.text = str(value)
                
                # 设置交替行颜色
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF7, 0xFA, 0xFC)
                
                # 设置字体
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(12)
        
        return table
    
    def _apply_chart_style(self, chart: Any):
        """应用图表样式"""
        # 设置图例
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        
        # 设置颜色
        plot = chart.plots[0]
        for i, series in enumerate(plot.series):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = self.colors[i % len(self.colors)]
        
        # 设置坐标轴
        if hasattr(chart, 'value_axis'):
            chart.value_axis.has_major_gridlines = True
            chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        
        if hasattr(chart, 'category_axis'):
            chart.category_axis.has_major_gridlines = False


class DataVisualizer:
    """数据可视化器"""
    
    def __init__(self):
        self.chart_generator = ChartGenerator()
    
    def create_visualization(self, slide: Any, chart_type: str, data: dict) -> Any:
        """创建数据可视化
        
        Args:
            slide: 幻灯片对象
            chart_type: 图表类型
            data: 数据
            
        Returns:
            图表对象
        """
        chart_generators = {
            "bar": self.chart_generator.generate_bar_chart,
            "line": self.chart_generator.generate_line_chart,
            "pie": self.chart_generator.generate_pie_chart,
            "table": self.chart_generator.generate_table,
        }
        
        generator = chart_generators.get(chart_type)
        if generator:
            return generator(slide, data)
        
        # 默认使用柱状图
        return self.chart_generator.generate_bar_chart(slide, data)
    
    def format_metric(self, value: Any, format_type: str = "number") -> str:
        """格式化指标值
        
        Args:
            value: 指标值
            format_type: 格式类型
            
        Returns:
            格式化后的字符串
        """
        if format_type == "percent":
            return f"{value:.1%}"
        elif format_type == "currency":
            return f"¥{value:,.2f}"
        elif format_type == "number":
            if isinstance(value, float):
                return f"{value:,.2f}"
            return f"{value:,}"
        elif format_type == "growth":
            if value > 0:
                return f"↑{value:.1%}"
            elif value < 0:
                return f"↓{abs(value):.1%}"
            else:
                return "持平"
        
        return str(value)
