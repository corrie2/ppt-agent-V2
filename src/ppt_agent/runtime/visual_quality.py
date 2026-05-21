from __future__ import annotations

import json
from pathlib import Path

import httpx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pydantic import BaseModel, Field, ValidationError

from ppt_agent.domain.models import PptSpec, SlideSpec
from ppt_agent.llm.planner import PlannerConfigError
from ppt_agent.runtime.agent_llm import AgentLlmConfig, generate_agent_json, load_agent_llm_config


class VisualQualityIssue(BaseModel):
    id: str
    severity: str = "warning"
    slide_no: int | None = None
    message: str
    suggested_fix: str | None = None
    rework_target: str | None = None


class SlideVisualScore(BaseModel):
    slide_no: int
    score: float
    layout: str = ""
    density: str = "standard"
    strengths: list[str] = Field(default_factory=list)
    issues: list[str] = Field(default_factory=list)


class VisualQualityReport(BaseModel):
    agent: str = "visual_quality_evaluator"
    model: str = "deepseek-v4-pro"
    ok: bool
    score: float
    summary: str
    slide_scores: list[SlideVisualScore] = Field(default_factory=list)
    issues: list[VisualQualityIssue] = Field(default_factory=list)
    metrics: dict = Field(default_factory=dict)
    rework_target: str | None = None


def evaluate_pptx_visual_quality(
    spec: PptSpec,
    pptx_path: Path,
    *,
    report_path: Path | None = None,
    config: AgentLlmConfig | None = None,
) -> VisualQualityReport:
    metrics = collect_pptx_visual_metrics(spec, pptx_path)
    deterministic = deterministic_visual_quality_report(spec, metrics)
    resolved_config = config or load_agent_llm_config()
    llm_report = _visual_quality_with_llm(
        spec,
        metrics,
        deterministic,
        context=_load_pipeline_context(spec),
        config=resolved_config,
    )
    report = llm_report or deterministic
    if report_path is not None:
        write_visual_quality_report(report, report_path)
    return report


def visual_quality_report_path(pptx_path: Path) -> Path:
    return pptx_path.with_suffix(".visual_quality_report.json")


def write_visual_quality_report(report: VisualQualityReport, path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(report.model_dump(mode="json"), ensure_ascii=False, indent=2), encoding="utf-8")


def collect_pptx_visual_metrics(spec: PptSpec, pptx_path: Path) -> dict:
    prs = Presentation(pptx_path)
    slides = []
    for index, slide in enumerate(prs.slides, start=1):
        text = []
        picture_count = 0
        table_count = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_count += 1
            if getattr(shape, "has_table", False):
                table_count += 1
            if getattr(shape, "has_text_frame", False):
                text.append(shape.text or "")
        text_body = "\n".join(part for part in text if part)
        spec_slide = spec.slides[index - 1] if index <= len(spec.slides) else None
        slides.append(
            {
                "slide_no": index,
                "shape_count": len(slide.shapes),
                "text_box_count": sum(1 for shape in slide.shapes if getattr(shape, "has_text_frame", False)),
                "text_chars": len(text_body),
                "picture_count": picture_count,
                "table_count": table_count,
                "planned_layout": (spec_slide.layout or spec_slide.layout_hint or "") if spec_slide else "",
                "planned_role": spec_slide.role if spec_slide else "",
                "planned_figure_ids": list(spec_slide.content.figure_ids) if spec_slide else [],
                "planned_bullet_count": len(spec_slide.bullets) if spec_slide else 0,
            }
        )
    return {
        "pptx_path": str(pptx_path),
        "slide_count": len(prs.slides),
        "planned_slide_count": len(spec.slides),
        "slides": slides,
    }


def deterministic_visual_quality_report(spec: PptSpec, metrics: dict) -> VisualQualityReport:
    issues: list[VisualQualityIssue] = []
    slide_scores: list[SlideVisualScore] = []
    if metrics["slide_count"] != metrics["planned_slide_count"]:
        issues.append(
            VisualQualityIssue(
                id="deck:slide_count_mismatch",
                severity="error",
                message="Rendered PPTX slide count differs from PptSpec slide count.",
                suggested_fix="Regenerate PptSpec or inspect PPTX build failures.",
                rework_target="page_generator",
            )
        )

    for slide_metric in metrics["slides"]:
        slide_no = slide_metric["slide_no"]
        score = 0.9
        slide_issues = []
        if slide_metric["shape_count"] <= 1:
            score -= 0.35
            slide_issues.append("slide may be blank or visually incomplete")
            issues.append(
                VisualQualityIssue(
                    id=f"slide-{slide_no:03d}:low_shape_count",
                    severity="error",
                    slide_no=slide_no,
                    message="Slide has very few rendered shapes and may be blank.",
                    suggested_fix="Check renderer layout fallback and picture/text placement.",
                    rework_target="page_generator",
                )
            )
        if slide_metric["text_chars"] > 950:
            score -= 0.2
            slide_issues.append("text density is high")
            issues.append(
                VisualQualityIssue(
                    id=f"slide-{slide_no:03d}:text_too_dense",
                    severity="warning",
                    slide_no=slide_no,
                    message="Slide has high rendered text density.",
                    suggested_fix="Reduce visible text or use a more visual layout.",
                    rework_target="page_designer",
                )
            )
        if slide_metric["planned_figure_ids"] and slide_metric["picture_count"] == 0:
            score -= 0.25
            slide_issues.append("planned figures were not rendered as pictures")
            issues.append(
                VisualQualityIssue(
                    id=f"slide-{slide_no:03d}:missing_planned_picture",
                    severity="warning",
                    slide_no=slide_no,
                    message="Slide planned figure evidence but rendered no picture shape.",
                    suggested_fix="Check evidence figure paths and renderer figure layout support.",
                    rework_target="renderer_engineer",
                )
            )
        slide_scores.append(
            SlideVisualScore(
                slide_no=slide_no,
                score=round(max(0.0, min(1.0, score)), 2),
                layout=slide_metric.get("planned_layout") or "",
                density="high" if slide_metric["text_chars"] > 950 else "standard",
                strengths=[] if slide_issues else ["basic visual structure is present"],
                issues=slide_issues,
            )
        )

    deck_score = round(sum(slide.score for slide in slide_scores) / max(1, len(slide_scores)), 2)
    error = any(issue.severity == "error" for issue in issues)
    rework_target = _first_rework_target(issues)
    return VisualQualityReport(
        ok=not error and deck_score >= 0.75,
        score=deck_score,
        summary="Deterministic visual quality checks completed from PPTX structure metrics.",
        slide_scores=slide_scores,
        issues=issues,
        metrics=metrics,
        rework_target=rework_target,
    )


def _visual_quality_with_llm(
    spec: PptSpec,
    metrics: dict,
    deterministic: VisualQualityReport,
    *,
    context: dict,
    config: AgentLlmConfig,
) -> VisualQualityReport | None:
    try:
        payload = generate_agent_json(
            "visual_quality_evaluator",
            system_prompt=_visual_quality_prompt(),
            user_payload={
                "ppt_spec_summary": _ppt_spec_summary(spec),
                "visual_metrics": metrics,
                "deterministic_report": deterministic.model_dump(mode="json"),
                "pipeline_context": context,
            },
            config=config,
        )
        if payload is None:
            return None
        return VisualQualityReport.model_validate(payload)
    except (httpx.HTTPError, PlannerConfigError, ValueError, ValidationError):
        if config.fallback_to_deterministic:
            return None
        raise


def _visual_quality_prompt() -> str:
    return (
        "Evaluate the final generated PPTX visual quality for a research or presentation deck. "
        "You receive structured PPTX metrics instead of screenshots. Return JSON only matching "
        "{agent, model, ok, score, summary, slide_scores, issues, metrics, rework_target}. "
        "Judge whether pages are too dense, visually incomplete, missing planned figures, or not suitable for a graduate presentation. "
        "Use rework_target page_designer for visual hierarchy/layout issues, renderer_engineer for renderer capability gaps, "
        "page_generator for mapping/build issues, or render_review for missed review issues. "
        "Do not rewrite slide content and do not claim to have seen screenshots."
    )


def _ppt_spec_summary(spec: PptSpec) -> dict:
    return {
        "title": spec.title,
        "audience": spec.audience,
        "slide_count": len(spec.slides),
        "slides": [
            {
                "slide_no": index,
                "id": slide.id,
                "role": slide.role,
                "title": slide.title,
                "layout": slide.layout or slide.layout_hint,
                "visual_type": slide.visual_type,
                "bullet_count": len(slide.bullets),
                "figure_ids": list(slide.content.figure_ids),
            }
            for index, slide in enumerate(spec.slides, start=1)
        ],
    }


def _load_pipeline_context(spec: PptSpec) -> dict:
    source = spec.source_digest or {}
    artifacts = source.get("artifacts") if isinstance(source, dict) else None
    if not isinstance(artifacts, dict):
        return {}
    context = {}
    for key in ("page_design", "renderer_engineer_report", "render_review_report"):
        path_value = artifacts.get(key)
        if not path_value:
            continue
        path = Path(path_value)
        if not path.exists():
            continue
        try:
            context[key] = json.loads(path.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError):
            continue
    return context


def _first_rework_target(issues: list[VisualQualityIssue]) -> str | None:
    for severity in ("error", "warning"):
        for issue in issues:
            if issue.severity == severity and issue.rework_target:
                return issue.rework_target
    return None
