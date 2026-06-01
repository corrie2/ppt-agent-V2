"""汇报PPT模板生成器"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# 汇报专用配色方案
EXECUTIVE_COLORS = {
    "primary": RGBColor(0x1A, 0x36, 0x5D),      # #1a365d 深蓝
    "secondary": RGBColor(0x2B, 0x6C, 0xB0),    # #2b6cb0 浅蓝
    "accent": RGBColor(0x31, 0x82, 0xCE),        # #3182ce 强调蓝
    "background": RGBColor(0xFF, 0xFF, 0xFF),    # #ffffff 白色
    "text_primary": RGBColor(0x1A, 0x20, 0x2C),  # #1a202c 深灰
    "text_secondary": RGBColor(0x4A, 0x55, 0x68),# #4a5568 中灰
    "border": RGBColor(0xE2, 0xE8, 0xF0),        # #e2e8f0 浅灰
    "success": RGBColor(0x38, 0xA1, 0x69),       # #38a169 绿色
    "warning": RGBColor(0xD6, 0x9E, 0x2E),       # #d69e2e 黄色
    "error": RGBColor(0xE5, 0x3E, 0x3E),         # #e53e3e 红色
}


class ExecutiveTemplateGenerator:
    """汇报PPT模板生成器"""
    
    def __init__(self, output_dir: str = "templates/executive-report"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.colors = EXECUTIVE_COLORS
    
    def generate_all_templates(self):
        """生成所有汇报模板"""
        templates = [
            ("quarterly-review", self._generate_quarterly_review),
            ("project-status", self._generate_project_status),
            ("annual-summary", self._generate_annual_summary),
            ("strategy-proposal", self._generate_strategy_proposal),
            ("data-analysis", self._generate_data_analysis),
        ]
        
        for name, generator in templates:
            print(f"生成模板: {name}")
            generator()
        
        print(f"所有模板已生成到: {self.output_dir}")
    
    def _create_presentation(self) -> Presentation:
        """创建16:9的演示文稿"""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        return prs
    
    def _add_title_slide(self, prs: Presentation, title: str, subtitle: str):
        """添加封面页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        
        # 添加背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors["primary"]
        
        # 添加标题
        left = Inches(1.5)
        top = Inches(2.5)
        width = Inches(10)
        height = Inches(1.5)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.colors["background"]
        p.alignment = PP_ALIGN.CENTER
        
        # 添加副标题
        top = Inches(4.2)
        height = Inches(1)
        
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = self.colors["border"]
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def _add_content_slide(self, prs: Presentation, title: str, bullets: list):
        """添加内容页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        
        # 添加标题栏
        left = Inches(0.8)
        top = Inches(0.4)
        width = Inches(11.5)
        height = Inches(0.8)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.colors["primary"]
        
        # 添加分隔线
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, Inches(1.3), width, Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors["secondary"]
        line.line.fill.background()
        
        # 添加内容
        top = Inches(1.8)
        height = Inches(5)
        
        content_box = slide.shapes.add_textbox(left, top, width, height)
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = f"• {bullet}"
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors["text_primary"]
            p.space_after = Pt(12)
        
        # 添加页码
        self._add_page_number(slide, prs)
        
        return slide
    
    def _add_chart_slide(self, prs: Presentation, title: str, chart_placeholder: str):
        """添加图表页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        
        # 添加标题
        left = Inches(0.8)
        top = Inches(0.4)
        width = Inches(11.5)
        height = Inches(0.8)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.colors["primary"]
        
        # 添加图表占位符
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(10)
        height = Inches(4.5)
        
        chart_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        chart_box.fill.solid()
        chart_box.fill.fore_color.rgb = self.colors["border"]
        chart_box.line.color.rgb = self.colors["text_secondary"]
        
        # 添加占位符文本
        text_frame = chart_box.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.text = chart_placeholder
        p.font.size = Pt(18)
        p.font.color.rgb = self.colors["text_secondary"]
        p.alignment = PP_ALIGN.CENTER
        
        # 添加页码
        self._add_page_number(slide, prs)
        
        return slide
    
    def _add_summary_slide(self, prs: Presentation, title: str, key_points: list, recommendations: list):
        """添加总结页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        
        # 添加标题
        left = Inches(0.8)
        top = Inches(0.4)
        width = Inches(11.5)
        height = Inches(0.8)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.colors["primary"]
        
        # 添加关键发现
        left = Inches(0.8)
        top = Inches(1.5)
        width = Inches(5.5)
        height = Inches(5)
        
        findings_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        findings_box.fill.solid()
        findings_box.fill.fore_color.rgb = self.colors["border"]
        findings_box.line.fill.background()
        
        text_frame = findings_box.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.text = "关键发现"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.colors["primary"]
        
        for point in key_points:
            p = text_frame.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors["text_primary"]
            p.space_after = Pt(8)
        
        # 添加建议
        left = Inches(6.8)
        
        rec_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        rec_box.fill.solid()
        rec_box.fill.fore_color.rgb = self.colors["primary"]
        rec_box.line.fill.background()
        
        text_frame = rec_box.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.text = "建议"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.colors["background"]
        
        for rec in recommendations:
            p = text_frame.add_paragraph()
            p.text = f"• {rec}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors["background"]
            p.space_after = Pt(8)
        
        # 添加页码
        self._add_page_number(slide, prs)
        
        return slide
    
    def _add_qa_slide(self, prs: Presentation):
        """添加Q&A页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        
        # 添加背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors["primary"]
        
        # 添加Q&A文本
        left = Inches(1.5)
        top = Inches(2.5)
        width = Inches(10)
        height = Inches(1.5)
        
        qa_box = slide.shapes.add_textbox(left, top, width, height)
        qa_frame = qa_box.text_frame
        qa_frame.word_wrap = True
        
        p = qa_frame.paragraphs[0]
        p.text = "Q & A"
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = self.colors["background"]
        p.alignment = PP_ALIGN.CENTER
        
        # 添加提示
        top = Inches(4.5)
        height = Inches(1)
        
        hint_box = slide.shapes.add_textbox(left, top, width, height)
        hint_frame = hint_box.text_frame
        hint_frame.word_wrap = True
        
        p = hint_frame.paragraphs[0]
        p.text = "感谢聆听，欢迎提问"
        p.font.size = Pt(24)
        p.font.color.rgb = self.colors["border"]
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def _add_page_number(self, slide: Presentation, prs: Presentation):
        """添加页码"""
        left = Inches(12)
        top = Inches(7)
        width = Inches(1)
        height = Inches(0.4)
        
        page_box = slide.shapes.add_textbox(left, top, width, height)
        page_frame = page_box.text_frame
        
        p = page_frame.paragraphs[0]
        p.text = str(len(prs.slides))
        p.font.size = Pt(12)
        p.font.color.rgb = self.colors["text_secondary"]
        p.alignment = PP_ALIGN.RIGHT
    
    def _generate_quarterly_review(self):
        """生成季度汇报模板"""
        prs = self._create_presentation()
        
        # 封面
        self._add_title_slide(prs, "2024年Q4工作汇报", "部门名称 | 汇报人 | 日期")
        
        # 目录
        self._add_content_slide(prs, "目录", [
            "一、关键指标概览",
            "二、重点工作成果",
            "三、问题与挑战",
            "四、下季度计划",
            "五、总结与建议"
        ])
        
        # 关键指标
        self._add_chart_slide(prs, "关键指标概览", "[在此插入KPI图表]")
        
        # 重点工作成果
        self._add_content_slide(prs, "重点工作成果", [
            "项目A：完成率95%，超预期",
            "项目B：按时交付，获得客户好评",
            "项目C：成本控制在预算内",
            "团队建设：新增3名核心成员"
        ])
        
        # 问题与挑战
        self._add_content_slide(prs, "问题与挑战", [
            "资源紧张：人员配置不足",
            "技术债务：需要重构核心模块",
            "市场竞争：竞品功能快速迭代",
            "客户需求：需求变更频繁"
        ])
        
        # 下季度计划
        self._add_content_slide(prs, "下季度计划", [
            "完成项目D的开发和上线",
            "优化现有系统性能",
            "扩展团队规模",
            "加强客户沟通"
        ])
        
        # 总结
        self._add_summary_slide(prs, "总结与建议", [
            "Q4整体完成率92%",
            "客户满意度提升15%",
            "团队效率提升20%"
        ], [
            "增加资源投入",
            "优化流程",
            "加强培训"
        ])
        
        # Q&A
        self._add_qa_slide(prs)
        
        output_path = self.output_dir / "quarterly-review.pptx"
        prs.save(str(output_path))
        print(f"已生成: {output_path}")
    
    def _generate_project_status(self):
        """生成项目进展汇报模板"""
        prs = self._create_presentation()
        
        # 封面
        self._add_title_slide(prs, "项目进展汇报", "项目名称 | 汇报人 | 日期")
        
        # 项目概览
        self._add_content_slide(prs, "项目概览", [
            "项目目标：...",
            "项目周期：2024年Q1-Q4",
            "项目团队：10人",
            "当前阶段：开发中"
        ])
        
        # 进度跟踪
        self._add_chart_slide(prs, "进度跟踪", "[在此插入甘特图或进度条]")
        
        # 风险与问题
        self._add_content_slide(prs, "风险与问题", [
            "风险1：技术方案待验证",
            "风险2：第三方依赖不稳定",
            "问题1：需求变更频繁",
            "问题2：测试资源不足"
        ])
        
        # 资源使用
        self._add_chart_slide(prs, "资源使用", "[在此插入资源使用图表]")
        
        # 里程碑
        self._add_content_slide(prs, "里程碑", [
            "里程碑1：需求评审 ✓",
            "里程碑2：设计评审 ✓",
            "里程碑3：开发完成 - 进行中",
            "里程碑4：测试完成 - 待开始"
        ])
        
        # 下一步计划
        self._add_content_slide(prs, "下一步计划", [
            "完成核心功能开发",
            "进行集成测试",
            "准备上线部署",
            "编写用户文档"
        ])
        
        # Q&A
        self._add_qa_slide(prs)
        
        output_path = self.output_dir / "project-status.pptx"
        prs.save(str(output_path))
        print(f"已生成: {output_path}")
    
    def _generate_annual_summary(self):
        """生成年度总结模板"""
        prs = self._create_presentation()
        
        # 封面
        self._add_title_slide(prs, "2024年度工作总结", "部门名称 | 汇报人 | 日期")
        
        # 年度概览
        self._add_content_slide(prs, "年度概览", [
            "全年营收：1.2亿，同比增长30%",
            "客户数量：新增50家，续约率95%",
            "团队规模：从20人扩展到35人",
            "产品迭代：发布4个大版本"
        ])
        
        # 主要成就
        self._add_content_slide(prs, "主要成就", [
            "完成A轮融资",
            "产品获得行业大奖",
            "建立战略合作伙伴关系",
            "团队获得最佳团队奖"
        ])
        
        # 关键数据
        self._add_chart_slide(prs, "关键数据", "[在此插入年度数据图表]")
        
        # 挑战与成长
        self._add_content_slide(prs, "挑战与成长", [
            "市场竞争加剧",
            "技术人才短缺",
            "客户需求多样化",
            "通过创新和合作克服挑战"
        ])
        
        # 团队贡献
        self._add_content_slide(prs, "团队贡献", [
            "技术团队：完成核心架构重构",
            "产品团队：推出3个创新功能",
            "销售团队：超额完成业绩目标",
            "运营团队：客户满意度提升20%"
        ])
        
        # 未来展望
        self._add_content_slide(prs, "未来展望", [
            "2025年目标：营收翻倍",
            "产品战略：AI赋能",
            "市场扩展：进入新行业",
            "团队建设：引进高端人才"
        ])
        
        # 致谢
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors["primary"]
        
        left = Inches(1.5)
        top = Inches(3)
        width = Inches(10)
        height = Inches(1.5)
        
        thanks_box = slide.shapes.add_textbox(left, top, width, height)
        thanks_frame = thanks_box.text_frame
        thanks_frame.word_wrap = True
        
        p = thanks_frame.paragraphs[0]
        p.text = "感谢团队的辛勤付出！"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.colors["background"]
        p.alignment = PP_ALIGN.CENTER
        
        output_path = self.output_dir / "annual-summary.pptx"
        prs.save(str(output_path))
        print(f"已生成: {output_path}")
    
    def _generate_strategy_proposal(self):
        """生成战略提案模板"""
        prs = self._create_presentation()
        
        # 封面
        self._add_title_slide(prs, "战略提案", "提案名称 | 汇报人 | 日期")
        
        # 执行摘要
        self._add_content_slide(prs, "执行摘要", [
            "核心主张：...",
            "预期收益：...",
            "所需资源：...",
            "实施周期：..."
        ])
        
        # 问题/机会
        self._add_content_slide(prs, "问题/机会", [
            "市场机会：...",
            "客户需求：...",
            "竞争态势：...",
            "技术趋势：..."
        ])
        
        # 解决方案
        self._add_content_slide(prs, "解决方案", [
            "方案概述：...",
            "核心优势：...",
            "差异化：...",
            "可行性：..."
        ])
        
        # 市场分析
        self._add_chart_slide(prs, "市场分析", "[在此插入市场分析图表]")
        
        # 财务预测
        self._add_chart_slide(prs, "财务预测", "[在此插入财务预测图表]")
        
        # 风险评估
        self._add_content_slide(prs, "风险评估", [
            "市场风险：...",
            "技术风险：...",
            "运营风险：...",
            "应对措施：..."
        ])
        
        # Q&A
        self._add_qa_slide(prs)
        
        output_path = self.output_dir / "strategy-proposal.pptx"
        prs.save(str(output_path))
        print(f"已生成: {output_path}")
    
    def _generate_data_analysis(self):
        """生成数据分析汇报模板"""
        prs = self._create_presentation()
        
        # 封面
        self._add_title_slide(prs, "数据分析报告", "分析主题 | 汇报人 | 日期")
        
        # 分析目标
        self._add_content_slide(prs, "分析目标", [
            "核心问题：...",
            "分析范围：...",
            "数据来源：...",
            "分析方法：..."
        ])
        
        # 分析方法
        self._add_content_slide(prs, "分析方法", [
            "数据采集：...",
            "数据清洗：...",
            "分析模型：...",
            "验证方法：..."
        ])
        
        # 关键发现
        self._add_chart_slide(prs, "关键发现", "[在此插入数据图表]")
        
        # 业务洞察
        self._add_content_slide(prs, "业务洞察", [
            "洞察1：...",
            "洞察2：...",
            "洞察3：...",
            "行动建议：..."
        ])
        
        # 建议
        self._add_content_slide(prs, "建议", [
            "短期行动：...",
            "中期规划：...",
            "长期战略：...",
            "资源需求：..."
        ])
        
        # Q&A
        self._add_qa_slide(prs)
        
        output_path = self.output_dir / "data-analysis.pptx"
        prs.save(str(output_path))
        print(f"已生成: {output_path}")


if __name__ == "__main__":
    generator = ExecutiveTemplateGenerator()
    generator.generate_all_templates()
